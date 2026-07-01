"""Parsers for modern Office / ODT formats."""
from __future__ import annotations

import logging
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

# (NS constant defined further down next to the PPTX helpers.)

log = logging.getLogger(__name__)


def _cap(chunks: list[str], max_chars: int) -> str:
    out = "\n".join(c for c in chunks if c)
    return out[:max_chars]


def parse_docx(path: Path, max_chars: int) -> str:
    try:
        from docx import Document  # type: ignore
    except ImportError:
        log.warning("python-docx not installed; skipping %s", path)
        return ""
    try:
        doc = Document(str(path))
    except Exception as exc:
        log.warning("docx open failed %s: %s", path, exc)
        return ""
    chunks: list[str] = []
    total = 0
    for p in doc.paragraphs:
        txt = (p.text or "").strip()
        if not txt:
            continue
        chunks.append(txt)
        total += len(txt)
        if total >= max_chars:
            break
    # Include the first table as a fallback if paragraphs were empty
    if total < 40:
        for tbl in doc.tables[:2]:
            for row in tbl.rows:
                row_txt = " | ".join((c.text or "").strip() for c in row.cells)
                if row_txt.strip():
                    chunks.append(row_txt)
                    total += len(row_txt)
                    if total >= max_chars:
                        break
            if total >= max_chars:
                break
    return _cap(chunks, max_chars)


_PPTX_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _parse_pptx_via_xml(path: Path, max_chars: int) -> str:
    """Fallback PPTX text extractor that reads slide XML directly.

    Used when python-pptx's high-level API trips over malformed relationship
    parts (seen with some Korean templates: ``'list' object has no
    attribute 'rId'``).  We only need ``<a:t>`` text runs to feed the LLM.
    """
    chunks: list[str] = []
    total = 0
    try:
        with zipfile.ZipFile(path) as zf:
            slide_names = sorted(
                n for n in zf.namelist()
                if n.startswith("ppt/slides/slide") and n.endswith(".xml")
            )
            for name in slide_names[:30]:
                try:
                    with zf.open(name) as f:
                        xml = f.read()
                except Exception:
                    continue
                try:
                    root = ET.fromstring(xml)
                except ET.ParseError:
                    continue
                for elem in root.iter(f"{_PPTX_A_NS}t"):
                    txt = (elem.text or "").strip()
                    if not txt:
                        continue
                    chunks.append(txt)
                    total += len(txt)
                    if total >= max_chars:
                        return _cap(chunks, max_chars)
    except (zipfile.BadZipFile, KeyError) as exc:
        log.warning("pptx xml fallback failed %s: %s", path, exc)
    return _cap(chunks, max_chars)


def parse_pptx(path: Path, max_chars: int) -> str:
    """Extract text from a PPTX.

    We learned the hard way that python-pptx's high-level slide walk
    crashes with ``'list' object has no attribute 'rId'`` on essentially
    every Korean deck saved by PowerPoint Korean / Hancom Office (the
    relationship parts use a list shape that python-pptx's iterator
    doesn't anticipate).  The raw-XML extractor is both faster and
    100 % reliable for our purposes (we only need front-matter text),
    so we use it as the primary path now.  python-pptx is kept only
    as a tertiary fallback for documents the XML reader couldn't open
    (e.g. password-protected or ODP-via-extension decks).
    """
    text = _parse_pptx_via_xml(path, max_chars)
    if text:
        return text

    # XML reader returned nothing — try python-pptx as a last resort.
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ""
    try:
        pres = Presentation(str(path))
    except Exception as exc:
        log.debug("pptx open failed %s: %s", path, exc)
        return ""

    chunks: list[str] = []
    total = 0
    try:
        for slide in pres.slides[:20]:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                try:
                    paragraphs = shape.text_frame.paragraphs
                except Exception:
                    continue
                for para in paragraphs:
                    try:
                        runs = para.runs
                    except Exception:
                        continue
                    for run in runs:
                        try:
                            txt = (run.text or "").strip()
                        except Exception:
                            continue
                        if not txt:
                            continue
                        chunks.append(txt)
                        total += len(txt)
                        if total >= max_chars:
                            return _cap(chunks, max_chars)
    except Exception as exc:
        log.debug("pptx high-level walk failed %s: %s", path, exc)
    return _cap(chunks, max_chars)


def parse_xlsx(path: Path, max_chars: int) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError:
        log.warning("openpyxl not installed; skipping %s", path)
        return ""
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        log.warning("xlsx open failed %s: %s", path, exc)
        return ""
    chunks: list[str] = []
    total = 0
    try:
        for name in wb.sheetnames[:6]:
            chunks.append(f"# Sheet: {name}")
            ws = wb[name]
            for row in ws.iter_rows(max_row=50, values_only=True):
                cells = [str(c) for c in row if c is not None]
                if not cells:
                    continue
                row_txt = " | ".join(cells)
                chunks.append(row_txt)
                total += len(row_txt)
                if total >= max_chars:
                    return _cap(chunks, max_chars)
        return _cap(chunks, max_chars)
    finally:
        try:
            wb.close()
        except Exception:
            pass


def parse_legacy_office(path: Path, max_chars: int) -> str:
    """Best-effort text scrape from legacy binary Office formats
    (``.doc`` / ``.ppt`` / ``.pps`` / ``.xls``).  We don't try to fully
    decode the formats — indexing only needs *some* text from the
    cover page or first sheet — so we open the OLE compound storage
    via ``olefile`` and pick out long printable runs from the streams
    most likely to carry user text.
    """
    try:
        import olefile  # type: ignore
    except ImportError:
        return ""
    try:
        if not olefile.isOleFile(str(path)):
            return ""
        ole = olefile.OleFileIO(str(path))
    except Exception as exc:
        log.debug("legacy ole open failed %s: %s", path, exc)
        return ""

    # Stream names that typically contain readable text per format.
    candidates = {
        ".doc": ("WordDocument", "1Table", "0Table"),
        ".ppt": ("PowerPoint Document", "Pictures"),
        ".pps": ("PowerPoint Document", "Pictures"),
        ".xls": ("Workbook", "Book"),
    }.get(path.suffix.lower(), ())
    chunks: list[str] = []
    total = 0
    try:
        for entry in ole.listdir(streams=True):
            if not entry:
                continue
            top = entry[0]
            if candidates and top not in candidates:
                continue
            try:
                with ole.openstream("/".join(entry)) as s:
                    raw = s.read(64 * 1024)
            except Exception:
                continue
            text = _scrape_printable_from_bytes(raw)
            if not text:
                continue
            chunks.append(text)
            total += len(text)
            if total >= max_chars:
                break
    finally:
        try:
            ole.close()
        except Exception:
            pass
    return _cap(chunks, max_chars)


_PRINTABLE_RX = __import__("re").compile(
    r"[\x20-\x7e가-힣ㄱ-ㆎ一-鿿\s·…\-\.,\(\)\[\]\{\}\"\'‘’“”!?:;#&@/]+"
)


def _scrape_printable_from_bytes(data: bytes) -> str:
    """Scrape readable Korean / English runs from a raw binary stream.

    Tries UTF-16-LE first (common for modern .doc/.ppt) then CP949 /
    EUC-KR (common for older Korean files), then Latin-1 as a last
    resort.  Filters out short noise runs to avoid OLE structure bytes.
    """
    encodings = ("utf-16-le", "cp949", "euc-kr", "latin-1")
    best = ""
    for enc in encodings:
        try:
            decoded = data.decode(enc, errors="ignore")
        except Exception:
            continue
        runs = _PRINTABLE_RX.findall(decoded)
        text = "\n".join(r.strip() for r in runs if len(r.strip()) >= 4)
        if len(text) > len(best):
            best = text
    return best


def parse_odf(path: Path, max_chars: int) -> str:
    """Extract text from any OpenDocument container (odt/ods/odp/...)."""
    try:
        with zipfile.ZipFile(path) as z:
            with z.open("content.xml") as f:
                xml = f.read()
    except (zipfile.BadZipFile, KeyError) as exc:
        log.warning("odf open failed %s: %s", path, exc)
        return ""
    try:
        root = ET.fromstring(xml)
    except ET.ParseError as exc:
        log.warning("odf xml parse failed %s: %s", path, exc)
        return ""
    texts: list[str] = []
    total = 0
    for elem in root.iter():
        if elem.text and elem.text.strip():
            texts.append(elem.text.strip())
            total += len(elem.text)
            if total >= max_chars:
                break
    return _cap(texts, max_chars)
